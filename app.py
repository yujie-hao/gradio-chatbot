import os
import gradio as gr
import openai
import pandas as pd
import numpy as np
from docx import Document
from openai import OpenAI
from pptx import Presentation
import PyPDF2
from sklearn.metrics.pairwise import cosine_similarity
from dotenv import load_dotenv
from pathlib import Path
from typing import List, Dict, Tuple

# Load environment variables
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
# openai.api_key = ""
openai_client = OpenAI(api_key=openai.api_key)


class HrDocumentChatbot:
    def __init__(self):
        self.document_store = []
        self.embeddings = []
        self.processed_files = set()
        self.region_mapping = {
            'global': 'Global',
            'eu': 'Europe',
            'us': 'United States'
        }
        self.conversation_history = []  # Track conversation history

    def ingest_documents(self, base_folder_path: str):
        """Process all HR documents in the structured folder"""
        base_folder = Path(base_folder_path)
        if not base_folder.exists():
            return f"Base folder not found: {base_folder_path}"

        new_files = 0

        # Process global documents (no region specific)
        global_folder = base_folder / "global"
        if global_folder.exists():
            new_files += self._process_folder(global_folder, region='global')

        # Process EU documents
        eu_folder = base_folder / "eu"
        if eu_folder.exists():
            new_files += self._process_folder(eu_folder, region='eu')

        # Process US documents
        us_folder = base_folder / "us"
        if us_folder.exists():
            new_files += self._process_folder(us_folder, region='us')

        return f"Processed {new_files} new HR documents (Total: {len(self.document_store)} chunks)"

    def _process_folder(self, folder: Path, region: str = None) -> int:
        """Process all supported files in a specific folder"""
        new_files = 0
        for file_path in folder.glob("*"):
            if file_path.name in self.processed_files:
                continue

            if file_path.suffix.lower() == '.docx':
                new_files += self._process_docx(file_path, region)
            elif file_path.suffix.lower() == '.xlsx':
                new_files += self._process_xlsx(file_path, region)
            elif file_path.suffix.lower() == '.pptx':
                new_files += self._process_pptx(file_path, region)
            elif file_path.suffix.lower() == '.pdf':
                new_files += self._process_pdf(file_path, region)

        return new_files

    def _process_docx(self, file_path: Path, region: str = None) -> int:
        """Process a DOCX file into chunks with embeddings"""
        try:
            doc = Document(file_path)
            full_text = [para.text for para in doc.paragraphs if para.text.strip()]
            text = '\n'.join(full_text)

            # Split into chunks
            chunk_size = 2000
            chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

            for chunk in chunks:
                self._add_to_store(
                    content=chunk,
                    source=file_path.name,
                    content_type="docx",
                    region=region
                )

            self.processed_files.add(file_path.name)
            return 1
        except Exception as e:
            print(f"Error processing DOCX {file_path}: {e}")
            return 0

    def _process_xlsx(self, file_path: Path, region: str = None) -> int:
        """Process an XLSX file into rows with embeddings"""
        try:
            df = pd.read_excel(file_path)

            # Create enhanced text for each row
            for _, row in df.iterrows():
                row_text = ', '.join([f"{col}: {row[col]}" for col in df.columns])
                self._add_to_store(
                    content=row_text,
                    source=file_path.name,
                    content_type="xlsx",
                    region=region
                )

            self.processed_files.add(file_path.name)
            return 1
        except Exception as e:
            print(f"Error processing XLSX {file_path}: {e}")
            return 0

    def _process_pptx(self, file_path: Path, region: str = None) -> int:
        """Process a PPTX file into slides with embeddings"""
        try:
            prs = Presentation(file_path)

            for i, slide in enumerate(prs.slides):
                slide_text = []
                # Add slide title if exists
                if slide.shapes.title and slide.shapes.title.text.strip():
                    slide_text.append(f"Slide {i + 1} Title: {slide.shapes.title.text}")

                # Add all text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and not shape == slide.shapes.title:
                        if shape.text.strip():
                            slide_text.append(shape.text)

                if slide_text:
                    self._add_to_store(
                        content='\n'.join(slide_text),
                        source=file_path.name,
                        content_type="pptx",
                        slide_number=i + 1,
                        region=region
                    )

            self.processed_files.add(file_path.name)
            return 1
        except Exception as e:
            print(f"Error processing PPTX {file_path}: {e}")
            return 0

    def _process_pdf(self, file_path: Path, region: str = None) -> int:
        """Process a PDF file into pages with embeddings"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)

                for i, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        self._add_to_store(
                            content=text,
                            source=file_path.name,
                            content_type="pdf",
                            page_number=i + 1,
                            region=region
                        )

            self.processed_files.add(file_path.name)
            return 1
        except Exception as e:
            print(f"Error processing PDF {file_path}: {e}")
            return 0

    def _add_to_store(self, content: str, source: str, content_type: str, **metadata):
        """Add document chunk to store with embedding"""
        # Enhance content with region info if available
        enhanced_content = content
        if 'region' in metadata and metadata['region']:
            region_display = self.region_mapping.get(metadata['region'], metadata['region'])
            enhanced_content = f"[Region: {region_display}]\n{content}"

        embedding = self._get_embedding(enhanced_content)
        self.document_store.append({
            "content": enhanced_content,
            "source": source,
            "type": content_type,
            "embedding": embedding,
            "metadata": metadata
        })
        self.embeddings.append(embedding)

    def _get_embedding(self, text: str) -> List[float]:
        """Get embedding for text"""
        response = openai_client.embeddings.create(
            input=text,
            model="text-embedding-ada-002"
        )
        return response.data[0].embedding

    def find_relevant_chunks(self, query: str, top_k: int = 20) -> List[Dict]:
        """Find most relevant document chunks across all files"""
        if not self.document_store:
            return []

        query_embedding = self._get_embedding(query)
        similarities = cosine_similarity([query_embedding], self.embeddings)[0]
        top_indices = np.argsort(similarities)[-top_k:][::-1]
        return [self.document_store[i] for i in top_indices]

    def generate_response(self, query: str, history: List[Tuple[str, str]] = None) -> str:
        """Generate response using relevant chunks from all documents and conversation history"""
        relevant_chunks = self.find_relevant_chunks(query)

        if not relevant_chunks:
            return "No relevant information found in HR documents"

        # Build context from relevant chunks
        context = "Relevant information from HR documents:\n\n"
        for chunk in relevant_chunks:
            source_info = f"From {chunk['source']} ({chunk['type'].upper()}"

            # Add format-specific metadata
            if chunk['type'] == 'pptx':
                source_info += f", Slide {chunk['metadata'].get('slide_number', 'N/A')}"
            elif chunk['type'] == 'pdf':
                source_info += f", Page {chunk['metadata'].get('page_number', 'N/A')}"

            if 'region' in chunk['metadata'] and chunk['metadata']['region']:
                region_display = self.region_mapping.get(chunk['metadata']['region'], chunk['metadata']['region'])
                source_info += f", Region: {region_display}"

            context += f"{source_info}):\n{chunk['content']}\n\n"

        # Create prompt for OpenAI with conversation history
        messages = [
            {"role": "system", "content":
                "You are an HR assistant for ChemNovus Incorporated. Answer questions using ONLY the provided documents. "
                "Pay special attention to region-specific information (Global, US, Europe). "
                "When information differs between regions, present it in a clear comparison format. "
                "For global policies, indicate they apply worldwide. "
                "Format responses with clear section headers and bullet points."}
        ]

        # Add conversation history if available
        if history:
            for user_msg, bot_msg in history:
                messages.append({"role": "user", "content": user_msg})
                messages.append({"role": "assistant", "content": bot_msg})

        # Add current context and question
        messages.append({"role": "user", "content": f"{context}\n\nQuestion: {query}"})

        # Get response from OpenAI
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            temperature=0.3,
            max_tokens=5000,
        )

        # Update conversation history
        self.conversation_history.append((query, response.choices[0].message.content))

        return response.choices[0].message.content


# Initialize chatbot
chatbot = HrDocumentChatbot()


def process_folder(folder_path):
    """Handle folder processing"""
    if not folder_path:
        return "No folder selected"
    return chatbot.ingest_documents(folder_path)


def respond(message: str, history: List[Tuple[str, str]]):
    """Handle chat responses with history"""
    if history is None:
        history = []
    return chatbot.generate_response(message, history)


# Auto-detect dataset folder
current_dir = Path(__file__).parent if "__file__" in locals() else Path.cwd()
default_folder = current_dir / "dataset"

with gr.Blocks(title="HR Document Assistant", theme="soft") as demo:
    gr.Markdown("# 📂 ChemNovus Incorporated HR Assistant")
    gr.Markdown("Ask questions about global, US, and European HR policies, benefits, and procedures")

    with gr.Row():
        folder_input = gr.Textbox(
            label="Dataset Folder Path",
            value=str(default_folder) if default_folder.exists() else "",
            placeholder="Path to folder containing dataset/global, dataset/eu, dataset/us"
        )
        process_btn = gr.Button("Process HR Documents")

    status_output = gr.Textbox(label="Processing Status")
    process_btn.click(process_folder, inputs=folder_input, outputs=status_output)

    chatbot_interface = gr.ChatInterface(
        fn=respond,
        examples=[
            "Compare parental leave policies between US and Europe.",
            "Give me info about parental leave in the US.",
            "Give me info about parental leave.",
            "What onboarding tasks are, for Manager?",
            "What onboarding tasks are required for new hires?",
            "Show me the performance evaluation process.",
            "Show me the performance management cycle.",
            "Tell me more about goal setting and performance tracking.",
            "Tell me about health insurance benefits?",
            "What are the differences in health insurance between regions?",
            "List all global PTO policies.",
            "Brief the Disciplinary Actions.",
            "What are the health insurance benefits in the US?",
            "How about in Europe?",
            "What's the difference in vacation days between regions?"
        ],
        submit_btn="Ask HR Question",
        retry_btn=None,
        undo_btn=None,
        clear_btn="New Conversation"
    )


    # Add clear conversation history button
    def clear_history():
        chatbot.conversation_history = []
        return None


    chatbot_interface.clear_btn.click(
        fn=clear_history,
        inputs=None,
        outputs=None,
        queue=False
    )

if __name__ == "__main__":
    print("Expected folder structure:")
    print("dataset/")
    print("├── global/  # Global policies")
    print("├── eu/      # European region policies")
    print("└── us/      # US region policies")
    demo.launch(share=True)